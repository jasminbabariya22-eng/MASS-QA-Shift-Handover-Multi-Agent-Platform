import os
import uuid
import logfire
from typing import List, Optional
from PIL import Image

from app.config import settings
from app.ingestion.models import DocumentElement
from app.ingestion.parsers.base import DocumentParser
from app.ingestion.vision.analyzer import VisionAnalyzer, GeminiVisionAnalyzer, StubVisionAnalyzer


class FallbackDocumentParser(DocumentParser):
    """
    Fallback Multimodal Parser using pdfplumber/pypdfium2 for PDFs and python-pptx for PPTX.
    Guarantees document ingestion succeeds even if Docling encounters unsupported structures.
    """
    def __init__(self, vision_analyzer: Optional[VisionAnalyzer] = None):
        self.vision_analyzer = vision_analyzer or (
            GeminiVisionAnalyzer() if settings.VISION_ENABLED else StubVisionAnalyzer()
        )
        self.artifact_dir = settings.ARTIFACT_DIR

    def parse_document(
        self,
        file_path: str,
        document_id: str,
        document_name: str
    ) -> List[DocumentElement]:
        ext = document_name.lower().rsplit(".", 1)[-1]
        if ext == "pptx":
            return self._parse_pptx(file_path, document_id, document_name)
        return self._parse_pdf(file_path, document_id, document_name)

    def _parse_pdf(self, file_path: str, document_id: str, document_name: str) -> List[DocumentElement]:
        elements: List[DocumentElement] = []
        doc_artifact_path = os.path.join(self.artifact_dir, document_id)
        os.makedirs(doc_artifact_path, exist_ok=True)

        try:
            import pdfplumber
            import pypdfium2 as pdfium

            pdfium_doc = pdfium.PdfDocument(file_path)

            with pdfplumber.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    page_no = page_idx + 1

                    # 1. Extract Page Text
                    text_content = page.extract_text() or ""
                    if text_content.strip():
                        elem = DocumentElement(
                            element_id=f"{document_id}-P{page_no:03d}-TXT",
                            document_id=document_id,
                            document_name=document_name,
                            file_type="pdf",
                            content_type="text",
                            page_number=page_no,
                            text=text_content.strip()
                        )
                        elements.append(elem)

                    # 2. Extract Tables
                    tables = page.extract_tables() or []
                    for tbl_idx, table_data in enumerate(tables):
                        if not table_data or len(table_data) < 2:
                            continue
                        headers = [str(c or "").strip() for c in table_data[0]]
                        rows = [[str(val or "").strip() for val in r] for r in table_data[1:]]

                        semantic_lines = ["Table Data:"]
                        for r in rows[:15]:
                            semantic_lines.append(", ".join([f"{h}: {v}" for h, v in zip(headers, r)]))
                        semantic_text = "\n".join(semantic_lines)

                        elem = DocumentElement(
                            element_id=f"{document_id}-P{page_no:03d}-TBL{tbl_idx+1}",
                            document_id=document_id,
                            document_name=document_name,
                            file_type="pdf",
                            content_type="table",
                            page_number=page_no,
                            text=semantic_text,
                            table_data={"headers": headers, "rows": rows}
                        )
                        elements.append(elem)

                    # 3. Handle Visual Page Content (if text is minimal or images present)
                    if len(text_content.strip()) < 150 or not text_content.strip():
                        img_path = os.path.join(doc_artifact_path, f"page_{page_no:03d}.png")
                        try:
                            pdfium_page = pdfium_doc[page_idx]
                            pil_img = pdfium_page.render(scale=1.5).to_pil()
                            pil_img.save(img_path)

                            summary = self.vision_analyzer.analyze_image(img_path)
                            elem = DocumentElement(
                                element_id=f"{document_id}-P{page_no:03d}-VIS",
                                document_id=document_id,
                                document_name=document_name,
                                file_type="pdf",
                                content_type="diagram",
                                page_number=page_no,
                                text=summary,
                                image_path=img_path
                            )
                            elements.append(elem)
                        except Exception as vis_err:
                            logfire.warning(f"Could not analyze page visual for {document_name} P{page_no}: {vis_err}")

        except Exception as e:
            logfire.error(f"Fallback PDF parsing failed for {document_name}: {e}")

        return elements

    def _parse_pptx(self, file_path: str, document_id: str, document_name: str) -> List[DocumentElement]:
        elements: List[DocumentElement] = []
        doc_artifact_path = os.path.join(self.artifact_dir, document_id)
        os.makedirs(doc_artifact_path, exist_ok=True)

        try:
            from pptx import Presentation
            prs = Presentation(file_path)

            for slide_idx, slide in enumerate(prs.slides):
                slide_no = slide_idx + 1
                slide_title = ""
                slide_texts = []

                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = slide.shapes.title.text.strip()

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            t = paragraph.text.strip()
                            if t and t != slide_title:
                                slide_texts.append(t)

                    # Extract slide tables
                    if shape.has_table:
                        table_data = shape.table
                        headers = [cell.text.strip() for cell in table_data.rows[0].cells]
                        rows = []
                        for row in table_data.rows[1:]:
                            rows.append([cell.text.strip() for cell in row.cells])

                        semantic_lines = [f"Slide {slide_no} Table ({slide_title}):"]
                        for r in rows:
                            semantic_lines.append(", ".join([f"{h}: {v}" for h, v in zip(headers, r)]))

                        elem = DocumentElement(
                            element_id=f"{document_id}-S{slide_no:03d}-TBL",
                            document_id=document_id,
                            document_name=document_name,
                            file_type="pptx",
                            content_type="table",
                            slide_number=slide_no,
                            section=slide_title,
                            text="\n".join(semantic_lines),
                            table_data={"headers": headers, "rows": rows}
                        )
                        elements.append(elem)

                    # Extract slide images
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            image_bytes = shape.image.blob
                            img_ext = shape.image.ext
                            img_path = os.path.join(doc_artifact_path, f"slide_{slide_no:03d}_img.{img_ext}")
                            with open(img_path, "wb") as f:
                                f.write(image_bytes)

                            summary = self.vision_analyzer.analyze_image(img_path)
                            elem = DocumentElement(
                                element_id=f"{document_id}-S{slide_no:03d}-IMG",
                                document_id=document_id,
                                document_name=document_name,
                                file_type="pptx",
                                content_type="chart",
                                slide_number=slide_no,
                                section=slide_title,
                                text=summary,
                                image_path=img_path
                            )
                            elements.append(elem)
                        except Exception as img_err:
                            logfire.warning(f"Failed to extract PPTX image shape: {img_err}")

                # Combine slide text
                if slide_texts or slide_title:
                    full_slide_text = f"Slide Title: {slide_title}\n" + "\n".join(slide_texts)
                    elem = DocumentElement(
                        element_id=f"{document_id}-S{slide_no:03d}-TXT",
                        document_id=document_id,
                        document_name=document_name,
                        file_type="pptx",
                        content_type="slide",
                        slide_number=slide_no,
                        section=slide_title,
                        text=full_slide_text
                    )
                    elements.append(elem)

        except Exception as e:
            logfire.error(f"Fallback PPTX parsing failed for {document_name}: {e}")

        return elements
